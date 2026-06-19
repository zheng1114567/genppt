"""Python-native PPTX renderer — consumes 100% of DesignConcept to produce diverse slides."""

from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from .style import DesignConcept
from .tools.validators import resolve_structure


# ── Canvas ──
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


class DeckRenderer:
    """Renders a full PPTX from DesignConcept + slide specs. Zero information loss."""

    def __init__(
        self,
        concept: DesignConcept,
        slides: list[dict[str, Any]],
        brand: dict[str, str] | None = None,
        images_map: dict[int, str] | None = None,
    ):
        self.concept = concept
        self.slides = slides
        self.images_map = images_map or {}
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H

        # Resolve brand colors (concept takes precedence)
        b = brand or {}
        bg_hex = concept.background_hex.lstrip("#")
        self.colors = {
            "bg": b.get("bg", bg_hex),
            "surface": b.get("surface", _lighten_hex(concept.background_hex, 0.08)),
            "panel": b.get("panel", _panel_hex(concept.background_hex)),
            "ink": b.get("ink", concept.primary_hex.lstrip("#")),
            "muted": b.get("muted", _muted_hex(concept.primary_hex)),
            "line": b.get("line", _line_hex(concept.background_hex)),
            "accent": b.get("accent", concept.accent_hex.lstrip("#")),
            "accent_secondary": b.get("accent_secondary", concept.accent_secondary_hex.lstrip("#")),
            "accent_soft": b.get("accentSoft", _soften_hex(concept.accent_hex)),
            "green": b.get("green", concept.semantic_colors.get("positive", "059669")),
            "negative": b.get("negative", concept.semantic_colors.get("negative", "DC2626")),
            "on_dark": b.get("onDark", "F9FAFB"),
            "on_dark_subtle": b.get("onDarkSubtle", "9CA3AF"),
            "dark_bg": b.get("darkBg", "111827"),
            "dark_surface": b.get("darkSurface", "1F2937"),
        }

        # Typography from concept
        self.font = concept.font_family
        self.title_max = concept.max_title_size_pt
        self.body_base = concept.base_size_pt
        self.scale = concept.type_scale_ratio

        # Layout
        self.is_rounded = concept.shape_style == "rounded"
        self.corner = Inches(0.15) if self.is_rounded else Inches(0.04)
        self.margin = self._resolve_margin()
        self.decoration = concept.decoration_level

        # Layout routing (all accent placement is inline in handlers)

    # ── Public API ──

    def render(self, output_path: Path) -> Path:
        """Generate PPTX and save to output_path."""
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)

        for slide_spec in self.slides:
            page_num = int(slide_spec.get("index", 1))
            self._render_slide(slide_spec, page_num)

        self.prs.save(str(output_path))
        return output_path

    # ── Slide rendering ──

    def _render_slide(self, s: dict[str, Any], page_num: int) -> None:
        # Use blank layout — find it by name instead of hardcoded index
        blank_layout = None
        for layout in self.prs.slide_layouts:
            if layout.name == "Blank" or layout.name == "空白":
                blank_layout = layout
                break
        if blank_layout is None:
            blank_layout = self.prs.slide_layouts[-1]  # fallback to last layout
        slide = self.prs.slides.add_slide(blank_layout)

        spec = s.get("design", {})
        ct = s.get("color_treatment", {})
        bg = ct.get("bg", spec.get("bg", "light"))
        is_dark = bg == "dark"
        structure = resolve_structure(spec)

        # ── Sanitise: clamp dangerous parameter combos before rendering ──
        self._sanitise_spec(spec, s, structure)

        # Per-slide shape language overrides global corner radius
        slide_shape = s.get("shape_language", "")
        _shape_corners = {
            "organic_soft": Inches(0.2),
            "geometric_strict": Inches(0.02),
            "mixed": Inches(0.1),
            "minimal_lines_only": Inches(0.02),
        }
        self._slide_corner = _shape_corners.get(slide_shape, self.corner)

        # ── Background ──
        if is_dark:
            self._set_bg(slide, self.colors["dark_bg"])
            ink = self.colors["on_dark"]
            ink_subtle = self.colors["on_dark_subtle"]
            surface = self.colors["dark_surface"]
        elif bg == "accent_split":
            self._set_bg(slide, self.colors["bg"])
            self._add_rect(slide, 0, 0, SLIDE_W * 0.4, SLIDE_H, fill=self.colors["dark_bg"])
            ink = self.colors["ink"]
            ink_subtle = self.colors["muted"]
            surface = self.colors["surface"]
        elif bg == "accent_wash":
            self._set_bg(slide, self.colors["accent_soft"])
            ink = self.colors["ink"]
            ink_subtle = self.colors["muted"]
            surface = self.colors["surface"]
        else:
            self._set_bg(slide, self.colors["bg"])
            ink = self.colors["ink"]
            ink_subtle = self.colors["muted"]
            surface = self.colors["surface"]

        # ── Route to layout handler ──
        handler = {
            "hero_cover": self._render_hero,
            "centered": self._render_centered,
            "title_top": self._render_title_top,
            "title_split": self._render_title_split,
            "title_visual": self._render_title_visual,
            "grid": self._render_grid,
            "vertical_stack": self._render_vertical_stack,
            "accent_panel": self._render_accent_panel,
            "fishbone": self._render_fishbone,
            "timeline": self._render_timeline,
            "kpi_cards": self._render_kpi_cards,
            "comparison_table": self._render_comparison_table,
            "process_flow": self._render_process_flow,
            "quote_callout": self._render_quote_callout,
            "swot": self._render_swot,
            "funnel": self._render_funnel,
            "matrix_2x2": self._render_matrix_2x2,
            "agenda": self._render_agenda,
            "closing_cta": self._render_closing_cta,
        }.get(structure, self._render_title_top)

        # Normalize body: LLM may output string instead of list
        body_raw = s.get("body")
        if isinstance(body_raw, str):
            s["body"] = [body_raw]
        elif body_raw is None:
            s["body"] = []

        handler(slide, s, spec, ct, ink, ink_subtle, surface, is_dark, page_num)

        # ── Decorative image for content slides (small, bottom-right, non-intrusive) ──
        page_idx = int(s.get("index", page_num))
        if page_idx > 1:
            # Skip if layout already handles its own visual area
            _img_aware_layouts = {"hero_cover", "title_split", "title_visual", "closing_cta"}
            if structure not in _img_aware_layouts:
                slide_img = self.images_map.get(page_idx)
                if slide_img and Path(slide_img).exists():
                    img_w = Inches(2.2)
                    img_h = Inches(1.5)
                    self._add_image(slide, slide_img,
                                   SLIDE_W - self.margin - img_w,
                                   SLIDE_H - self.margin - img_h - Inches(0.1),
                                   img_w, img_h)

        # ── Page number on all slides ──
        if page_num > 1:  # skip cover
            self._add_text(slide, str(page_num),
                           SLIDE_W - Inches(0.8), SLIDE_H - Inches(0.45),
                           Inches(0.6), Inches(0.3),
                           size=Pt(8), color=ink_subtle, align="right")

        # ── Chart (only on layouts that support it) ──
        _chart_safe_layouts = {"title_top", "title_split", "title_visual", "grid", "centered", "vertical_stack"}
        if s.get("chart_spec") and structure in _chart_safe_layouts:
            self._render_chart(slide, s, spec, ct, ink)

    # ── Layout handlers ──

    # ── Shared layout primitives (all spacing flows through _gap) ──

    # ── Layout safeguard ──

    # Layouts that use card/table/grid structures with fixed space budgets.
    # These MUST NOT receive airy spacing or large title_size from the LLM.
    _CARD_LAYOUTS = {"kpi_cards", "swot", "funnel", "matrix_2x2", "comparison_table", "agenda"}
    _FLOW_LAYOUTS = {"process_flow", "timeline", "fishbone"}
    _STACK_LAYOUTS = {"vertical_stack", "grid"}

    def _sanitise_spec(self, spec: dict[str, Any], slide: dict[str, Any], structure: str) -> None:
        """Clamp dangerous LLM parameter choices before they hit layout handlers.

        The Design Agent sometimes picks title_size=42 + spacing=airy for a
        kpi_cards slide, which overflows the bottom margin.  This is a
        renderer-side safety net — the prompt already discourages these combos,
        but if the LLM ignores it, we auto-correct here.
        """
        title_size = spec.get("title_size", 28)
        spacing = spec.get("spacing", "normal")
        body = slide.get("body") or []
        body_len = len(body) if isinstance(body, list) else 1

        # ── Card layouts: no airy, moderate title ──
        if structure in self._CARD_LAYOUTS:
            if spacing == "airy":
                spec["spacing"] = "normal"
            if body_len >= 4 and spacing == "normal":
                spec["spacing"] = "compact"
            if title_size > 32:
                spec["title_size"] = 32

        # ── Flow layouts: no airy ──
        if structure in self._FLOW_LAYOUTS:
            if spacing == "airy":
                spec["spacing"] = "normal"
            if title_size > 34:
                spec["title_size"] = 34

        # ── Stack layouts: compact when content-heavy ──
        if structure in self._STACK_LAYOUTS:
            if body_len >= 4 and spacing != "compact":
                spec["spacing"] = "compact"
            if title_size > 30:
                spec["title_size"] = 30

        # ── Hero / closing / quote: allow airy, allow large title (no clamp) ──

    # ── Margin resolution ──

    _SPACING_MUL = {"tight": 0.55, "normal": 1.0, "airy": 1.6}

    def _gap(self, spec: dict, base: float):
        """Spacing-aware gap. Returns Inches.

        All layout handlers MUST use this for every internal gap — it is the
        single point where per-page spacing (tight/normal/airy) takes effect.
        """
        mul = self._SPACING_MUL.get(spec.get("spacing", "normal"), 1.0)
        return Inches(base * mul)

    def _draw_accent(self, slide, ct, x, y, h) -> float:
        """Draw accent marker. Returns how much vertical space it consumed."""
        placement = ct.get("accent_placement", "left_bar")
        if placement == "left_bar":
            self._add_rect(slide, x, y, Inches(0.08), h, fill=self.colors["accent"])
        elif placement == "top_strip":
            self._add_rect(slide, x, y, Inches(2.5), Inches(0.04), fill=self.colors["accent"])
        elif placement == "spot":
            self._add_circle(slide, x + Inches(0.04), y, Inches(0.22), fill=self.colors["accent"])
        elif placement == "text_highlight":
            return 0.0  # handled at text level, no shape drawn
        # "none" also returns 0
        return 0.0

    def _draw_header(self, slide, headline, kicker, spec, ct, ink, ink_subtle, m) -> float:
        """Draw title zone (accent + optional kicker + headline + separator).

        Returns the y position just below the separator — where body content should start.
        """
        title_size = spec.get("title_size", 28)
        spacing = spec.get("spacing", "normal")

        y = Inches(0.25)
        accent_h = self._gap(spec, 0.6)
        self._draw_accent(slide, ct, m, y, accent_h)

        header_y = y + accent_h + self._gap(spec, 0.1)
        if kicker:
            kicker_h = self._gap(spec, 0.3)
            self._add_text(slide, kicker.upper(), m, header_y, Inches(5), kicker_h,
                           size=Pt(9), color=ink_subtle, bold=True)
            header_y += kicker_h + self._gap(spec, 0.05)

        headline_h = self._gap(spec, 0.7)
        self._add_text(slide, headline, m, header_y, SLIDE_W - m * 2, headline_h,
                       size=Pt(title_size), color=ink, bold=True)

        sep_y = header_y + headline_h + self._gap(spec, 0.1)
        self._add_rect(slide, m, sep_y, Inches(2.5), Inches(0.02), fill=self.colors["line"])

        return sep_y + self._gap(spec, 0.25)

    # ── Individual layout handlers ──

    def _render_hero(self, slide, s, spec, ct, ink, ink_subtle, surface, is_dark, page_num):
        headline = str(s.get("headline", ""))
        kicker = str(s.get("kicker", s.get("intent", "")))
        body = s.get("body") or []
        title_size = spec.get("title_size", 48)
        m = self.margin
        page_idx = int(s.get("index", page_num))

        # Full-bleed dark background for cover impact
        self._set_bg(slide, self.colors["dark_bg"])
        dark_ink = self.colors["on_dark"]
        dark_subtle = self.colors["on_dark_subtle"]

        # Refined left content area — no full-height color block
        content_x = Inches(1.2)
        content_w = Inches(5.5)

        # Image on the right — if available
        hero_img = self.images_map.get(page_idx)
        if hero_img and Path(hero_img).exists():
            self._add_image(slide, hero_img, Inches(7.2), 0, Inches(6.1), SLIDE_H)
            # Subtle gradient overlay on image for text readability
            self._add_rect(slide, Inches(7.2), 0, Inches(6.1), SLIDE_H,
                           fill=self.colors["dark_bg"])
        else:
            # Refined right-side decoration: thin accent rules + subtle geometry
            # Outer frame
            self._add_rect(slide, Inches(7.5), Inches(1.2), Inches(4.5), Inches(5.1),
                           fill=None, line=self.colors["accent"])
            # Inner offset frame
            self._add_rect(slide, Inches(7.9), Inches(1.55), Inches(3.7), Inches(4.4),
                           fill=None, line=self.colors["accent_secondary"])
            # Accent corner mark — top-right of outer frame
            self._add_rect(slide, Inches(11.55), Inches(1.2), Inches(0.45), Inches(0.06),
                           fill=self.colors["accent"])
            self._add_rect(slide, Inches(11.55), Inches(1.2), Inches(0.06), Inches(0.45),
                           fill=self.colors["accent"])
            # Accent corner mark — bottom-left of outer frame
            self._add_rect(slide, Inches(7.5), Inches(5.85), Inches(0.45), Inches(0.06),
                           fill=self.colors["accent_secondary"])
            self._add_rect(slide, Inches(7.5), Inches(5.85), Inches(0.06), Inches(0.45),
                           fill=self.colors["accent_secondary"])
            # Small decorative circle at center of right area
            dot_r = Inches(0.15)
            self._add_circle(slide, Inches(9.5) - dot_r, Inches(3.45) - dot_r,
                             dot_r * 2, fill=self.colors["accent"])

        # Thin accent rule at top of content area
        self._add_rect(slide, content_x, Inches(1.5), Inches(1.8), Inches(0.045),
                       fill=self.colors["accent"])

        # Headline
        hl_size = Pt(min(max(title_size, 38), 48))
        self._add_text(slide, headline,
                       content_x, Inches(1.9),
                       content_w, Inches(3.2),
                       size=hl_size, color=dark_ink, bold=True)

        # Thin separator between headline and subtitle
        self._add_rect(slide, content_x, Inches(5.3), Inches(1.2), Inches(0.03),
                       fill=self.colors["accent_secondary"])

        # Kicker / subtitle
        if kicker:
            self._add_text(slide, kicker, content_x, Inches(5.6), content_w, Inches(0.5),
                           size=Pt(14), color=dark_subtle)

        # Body points
        if body:
            for j, line in enumerate(body[:2]):
                self._add_text(slide, f"▸ {str(line)}",
                               content_x, Inches(6.3 + j * 0.45),
                               content_w, Inches(0.45),
                               size=Pt(11), color=dark_subtle)

    def _render_centered(self, slide, s, spec, ct, ink, ink_subtle, surface, is_dark, page_num):
        headline = str(s.get("headline", ""))
        body = s.get("body") or []
        title_size = spec.get("title_size", 38)
        m = Inches(1.5)
        w = SLIDE_W - m * 2

        # Accent bar centered
        bar_y = Inches(2.3)
        self._add_rect(slide, (SLIDE_W - Inches(3.5)) / 2, bar_y, Inches(3.5), Inches(0.05),
                       fill=self.colors["accent"])

        headline_h = self._gap(spec, 1.2)
        self._add_text(slide, headline, m, bar_y + self._gap(spec, 0.3), w, headline_h,
                       size=Pt(min(title_size + 4, 48)), color=ink, bold=True, align="center")

        if body:
            body_y = bar_y + self._gap(spec, 0.3) + headline_h + self._gap(spec, 0.5)
            self._add_body_text(slide, body, Inches(2.5), body_y, SLIDE_W - Inches(5),
                                SLIDE_H - body_y - self._gap(spec, 1.0),
                                size=Pt(self.body_base + 1), color=ink_subtle)

    def _render_title_top(self, slide, s, spec, ct, ink, ink_subtle, surface, is_dark, page_num):
        headline = str(s.get("headline", ""))
        kicker = str(s.get("kicker", ""))
        body = s.get("body") or []
        title_size = spec.get("title_size", 28)
        body_size = spec.get("body_size", self.body_base)
        body_cols = spec.get("body_columns", 1)
        m = self._resolve_margin()

        body_y = self._draw_header(slide, headline, kicker, spec, ct, ink, ink_subtle, m)
        body_h = SLIDE_H - body_y - self._gap(spec, 0.4)
        body_w = SLIDE_W - m * 2

        if body_cols >= 3 and len(body) >= 3:
            self._render_body_cards(slide, body, m, body_y, body_w, body_h, 3,
                                    body_size, ink, ink_subtle, surface)
        elif body_cols == 2 and len(body) >= 2:
            self._render_body_cards(slide, body, m, body_y, body_w, body_h, 2,
                                    body_size, ink, ink_subtle, surface)
        else:
            self._add_body_text(slide, body, m, body_y, body_w, body_h,
                                size=Pt(body_size), color=ink_subtle)

    def _render_title_split(self, slide, s, spec, ct, ink, ink_subtle, surface, is_dark, page_num):
        headline = str(s.get("headline", ""))
        body = s.get("body") or []
        title_size = spec.get("title_size", 28)
        body_size = spec.get("body_size", self.body_base)

        # ── Layout parameters (LLM-controlled, sensible defaults) ──
        proportions = spec.get("proportions") or {}
        visual_ratio = proportions.get("visual", 0.35)           # 视觉区占比 0.25~0.50
        visual_side = spec.get("visual_side", "right")           # "left" | "right"
        spacing = spec.get("spacing", "normal")                  # "tight" | "normal" | "airy"

        # Spacing multiplier: affects all internal gaps
        gap_m = {"tight": 0.55, "normal": 1.0, "airy": 1.6}.get(spacing, 1.0)
        m = self.margin

        # ── Header zone ──
        accent_h = Inches(0.6 * gap_m)
        header_gap = Inches(0.3 * gap_m)
        headline_h = Inches(0.6 * gap_m)
        header_total = Inches(0.25) + accent_h + header_gap + headline_h
        split_y = header_total + Inches(0.15 * gap_m)
        split_h = SLIDE_H - split_y - Inches(0.35 * gap_m)

        # Accent
        accent_placement = ct.get("accent_placement", "left_bar")
        if accent_placement == "left_bar":
            self._add_rect(slide, m, Inches(0.25), Inches(0.08), accent_h, fill=self.colors["accent"])
        elif accent_placement == "top_strip":
            self._add_rect(slide, m, Inches(0.2), Inches(2.5), Inches(0.04), fill=self.colors["accent"])
        elif accent_placement == "spot":
            self._add_circle(slide, m + Inches(0.06), Inches(0.25), Inches(0.22), fill=self.colors["accent"])

        # Headline
        headline_y = Inches(0.25) + accent_h + header_gap
        self._add_text(slide, headline, m, headline_y, SLIDE_W - m * 2, headline_h,
                       size=Pt(title_size), color=ink, bold=True)

        # ── Body + Visual split ──
        total_w = SLIDE_W - m * 2
        gap = Inches(0.25 * gap_m)
        visual_w = total_w * visual_ratio
        body_w = total_w - visual_w - gap

        # Position body and visual based on side preference
        if visual_side == "left":
            vx = m
            bx = m + visual_w + gap
        else:
            bx = m
            vx = m + body_w + gap

        # Body panel
        self._add_rect(slide, bx, split_y, body_w, split_h,
                       fill=surface, line=self.colors["line"], radius=self.corner)
        if body:
            body_pad = Inches(0.3 * gap_m)
            self._add_body_text(slide, body, bx + body_pad, split_y + body_pad,
                                body_w - body_pad * 2, split_h - body_pad * 2,
                                size=Pt(body_size), color=ink_subtle)

        # Visual panel
        page_idx = int(s.get("index", page_num))
        visual_img = self.images_map.get(page_idx)

        if visual_img and Path(visual_img).exists():
            self._add_image(slide, visual_img, vx, split_y, visual_w, split_h)
        else:
            self._add_rect(slide, vx, split_y, visual_w, split_h,
                           fill=self.colors["accent_soft"], line=self.colors["line"],
                           radius=self._slide_corner)
            # Subtle decoration: thin rule + dot
            rule_w = min(visual_w * 0.45, Inches(1.8))
            self._add_rect(slide, vx + (visual_w - rule_w) / 2,
                           split_y + split_h * 0.42,
                           rule_w, Inches(0.025), fill=self.colors["accent"])
            dot_r = Inches(0.07)
            self._add_circle(slide, vx + visual_w / 2 - dot_r,
                             split_y + split_h * 0.55 - dot_r,
                             dot_r * 2, fill=self.colors["accent_secondary"])

    def _render_title_visual(self, slide, s, spec, ct, ink, ink_subtle, surface, is_dark, page_num):
        headline = str(s.get("headline", ""))
        body = s.get("body") or []
        title_size = spec.get("title_size", 28)
        m = self.margin
        spacing = spec.get("spacing", "normal")
        gap_m = {"tight": 0.55, "normal": 1.0, "airy": 1.6}.get(spacing, 1.0)
        visual_side = spec.get("visual_side", "right")

        # Accent
        accent_placement = ct.get("accent_placement", "left_bar")
        if accent_placement == "left_bar":
            self._add_rect(slide, m, Inches(0.55), Inches(0.06), Inches(0.45 * gap_m), fill=self.colors["accent"])
        elif accent_placement == "top_strip":
            self._add_rect(slide, m, Inches(0.3), Inches(2.0), Inches(0.04), fill=self.colors["accent"])

        # Headline
        self._add_text(slide, headline, m, Inches(0.55), SLIDE_W - m * 2, Inches(0.6 * gap_m),
                       size=Pt(title_size), color=ink, bold=True)

        visual_y = Inches(1.4 * gap_m)
        visual_h = SLIDE_H - visual_y - Inches(0.4 * gap_m)
        total_w = SLIDE_W - m * 2
        gap = Inches(0.25 * gap_m)

        if body:
            body_w = Inches(3.5 * gap_m)
            visual_w = total_w - body_w - gap
            if visual_side == "left":
                vx, bx = m, m + visual_w + gap
            else:
                bx, vx = m, m + body_w + gap
            self._add_body_text(slide, body, bx, visual_y, body_w, visual_h,
                                size=Pt(self.body_base), color=ink_subtle)
        else:
            vx = m + (total_w - total_w * 0.7) / 2
            visual_w = total_w * 0.7

        self._add_rect(slide, vx, visual_y, visual_w, visual_h,
                       fill=self.colors["accent_soft"], line=self.colors["line"], radius=self._slide_corner)
        # Subtle central decoration
        rule_w = min(visual_w * 0.35, Inches(1.5))
        self._add_rect(slide, vx + (visual_w - rule_w) / 2,
                       visual_y + visual_h * 0.48,
                       rule_w, Inches(0.025), fill=self.colors["accent"])
        dot_r = Inches(0.06)
        self._add_circle(slide, vx + visual_w / 2 - dot_r,
                         visual_y + visual_h * 0.58 - dot_r,
                         dot_r * 2, fill=self.colors["accent_secondary"])

    def _render_grid(self, slide, s, spec, ct, ink, ink_subtle, surface, is_dark, page_num):
        headline = str(s.get("headline", ""))
        body = s.get("body") or []
        body_size = spec.get("body_size", self.body_base)
        m = self.margin
        has_chart = bool(s.get("chart_spec"))
        body_cols = spec.get("body_columns", min(3, len(body)) if body else 1)

        body_y = self._draw_header(slide, headline, "", spec, ct, ink, ink_subtle, m)
        body_bottom = SLIDE_H - Inches(3.0) if has_chart else SLIDE_H - self._gap(spec, 1.0)
        body_h = body_bottom - body_y

        if body_cols >= 2 and len(body) >= 2:
            self._render_body_cards(slide, body, m, body_y, SLIDE_W - m * 2, body_h,
                                    body_cols, body_size, ink, ink_subtle, surface)
        elif body:
            self._add_body_text(slide, body, m, body_y, SLIDE_W - m * 2, body_h,
                                size=Pt(body_size), color=ink_subtle)

    def _render_vertical_stack(self, slide, s, spec, ct, ink, ink_subtle, surface, is_dark, page_num):
        headline = str(s.get("headline", ""))
        body = s.get("body") or []
        body_size = spec.get("body_size", self.body_base)
        m = self.margin

        body_start = self._draw_header(slide, headline, "", spec, ct, ink, ink_subtle, m)
        n = min(len(body) or 3, 5)
        step_gap = self._gap(spec, 0.1)
        bottom_pad = self._gap(spec, 0.4)
        step_h = (SLIDE_H - body_start - bottom_pad - step_gap * (n - 1)) / n
        colors = [self.colors["accent"], self.colors["accent_secondary"],
                  self.colors["green"], self.colors["accent"],
                  self.colors["accent_secondary"], self.colors["green"]]

        for i in range(n):
            sy = body_start + (step_h + step_gap) * i
            clr = colors[i % len(colors)]
            self._add_rect(slide, m, sy, SLIDE_W - m * 2, step_h,
                           fill=surface if i % 2 == 0 else self.colors["bg"],
                           line=self.colors["line"], radius=self._slide_corner)
            self._add_rect(slide, m, sy, Inches(0.06), step_h, fill=clr)
            num_r = Inches(0.28)
            self._add_text(slide, str(i + 1), m + Inches(0.25), sy + step_h / 2 - num_r,
                           Inches(0.4), num_r * 2, size=Pt(16), color=clr, bold=True)
            if i < len(body):
                self._add_text(slide, str(body[i]), m + Inches(0.8), sy + self._gap(spec, 0.15),
                               SLIDE_W - m * 2 - Inches(1.1), step_h - self._gap(spec, 0.3),
                               size=Pt(body_size), color=ink)

    def _render_accent_panel(self, slide, s, spec, ct, ink, ink_subtle, surface, is_dark, page_num):
        """Asymmetric layout: dark left panel (~35%) + content right."""
        headline = str(s.get("headline", ""))
        body = s.get("body") or []
        title_size = spec.get("title_size", 28)
        body_size = spec.get("body_size", self.body_base)
        proportions = spec.get("proportions", {})
        panel_ratio = proportions.get("visual", 0.38)
        m = self.margin

        panel_w = SLIDE_W * panel_ratio
        content_x = panel_w + Inches(0.5)
        content_w = SLIDE_W - content_x - m

        # Dark left panel
        self._add_rect(slide, 0, 0, panel_w, SLIDE_H, fill=self.colors["dark_bg"])

        # Geometric accents on panel — sized relative to panel
        accent_s = Inches(min(1.6, panel_w.inches * 0.35))
        self._add_rect(slide, panel_w * 0.25, Inches(2.2), accent_s, accent_s,
                       fill=self.colors["accent"], radius=self.corner)
        small_s = Inches(accent_s.inches * 0.5)
        self._add_rect(slide, panel_w * 0.25 + Inches(0.4), Inches(4.0), small_s, small_s,
                       fill=self.colors["accent_secondary"], radius=self.corner)

        # Headline
        hl_y = self._gap(spec, 1.5)
        hl_h = self._gap(spec, 1.0)
        self._add_text(slide, headline, content_x, hl_y, content_w, hl_h,
                       size=Pt(min(title_size + 6, 44)), color=ink, bold=True)

        # Accent bar
        self._add_rect(slide, content_x, hl_y + hl_h + self._gap(spec, 0.1),
                       Inches(2.5), Inches(0.04), fill=self.colors["accent"])

        # Body
        if body:
            body_y = hl_y + hl_h + self._gap(spec, 0.5)
            body_h = SLIDE_H - body_y - self._gap(spec, 0.6)
            self._add_body_text(slide, body, content_x, body_y, content_w, body_h,
                                size=Pt(body_size), color=ink_subtle)

        # Bottom rule
        self._add_rect(slide, content_x, SLIDE_H - self._gap(spec, 0.5),
                       content_w, Inches(0.015), fill=self.colors["line"])

    # ── Rich layout: Fishbone (Ishikawa) ──

    def _render_fishbone(self, slide, s, spec, ct, ink, ink_subtle, surface, is_dark, page_num):
        """Fishbone / Ishikawa diagram for root-cause analysis."""
        headline = str(s.get("headline", ""))
        body = s.get("body") or []
        title_size = spec.get("title_size", 24)
        m = self.margin

        self._add_text(slide, headline, m, Inches(0.25), SLIDE_W - m * 2, Inches(0.45),
                       size=Pt(title_size), color=ink, bold=True)

        # Main spine
        spine_y = SLIDE_H * 0.48
        spine_x1 = Inches(1.8)
        spine_x2 = Inches(10.2)
        self._add_rect(slide, spine_x1, spine_y, spine_x2 - spine_x1, Inches(0.04),
                       fill=self.colors["accent"])

        # Head box (problem statement)
        head_text = str(s.get("kicker", "")) or "核心问题"
        head_w = Inches(2.2); head_h = Inches(0.7)
        head_x = spine_x2 + Inches(0.15)
        head_y = spine_y - head_h / 2
        self._add_rect(slide, head_x, head_y, head_w, head_h,
                       fill=self.colors["dark_bg"], radius=self._slide_corner)
        self._add_text(slide, head_text, head_x + Inches(0.15), head_y + Inches(0.1),
                       head_w - Inches(0.3), head_h - Inches(0.2),
                       size=Pt(11), color=self.colors["on_dark"], bold=True, align="center")

        # Ribs
        causes = body if body else []
        rib_colors = [self.colors["accent"], self.colors["accent_secondary"],
                      self.colors["green"], "#8B5CF6", self.colors["accent"],
                      self.colors["green"]]
        for i, cause in enumerate(causes[:6]):
            text = str(cause)
            is_above = i % 2 == 0
            t = i / max(len(causes[:6]) - 1, 1) if len(causes) > 1 else 0.5
            rib_x = spine_x1 + Inches(1.2) + (spine_x2 - spine_x1 - Inches(2.4)) * t
            rib_len = Inches(1.5)
            clr = rib_colors[i % len(rib_colors)]

            # Draw rib as a thin rectangle (more reliable than connector)
            tip_y = spine_y - rib_len if is_above else spine_y + rib_len
            if is_above:
                self._add_rect(slide, rib_x - Inches(0.01), tip_y, Inches(0.02), rib_len,
                               fill=clr)
            else:
                self._add_rect(slide, rib_x - Inches(0.01), spine_y, Inches(0.02), rib_len,
                               fill=clr)
            dot_r = Inches(0.08)
            self._add_circle(slide, rib_x - dot_r, tip_y - dot_r, dot_r * 2, fill=clr)

            label_y = tip_y - Inches(0.55) if is_above else tip_y + Inches(0.15)
            self._add_text(slide, text, rib_x - Inches(1.1), label_y,
                           Inches(2.2), Inches(0.5),
                           size=Pt(8), color=ink_subtle, align="center")

    # ── Rich layout: Timeline ──

    def _render_timeline(self, slide, s, spec, ct, ink, ink_subtle, surface, is_dark, page_num):
        """Horizontal timeline with milestone nodes."""
        headline = str(s.get("headline", ""))
        body = s.get("body") or []
        m = self.margin

        body_start = self._draw_header(slide, headline, "", spec, ct, ink, ink_subtle, m)
        tl_y = body_start + self._gap(spec, 1.0)
        self._add_rect(slide, m + Inches(0.8), tl_y,
                       SLIDE_W - m * 2 - Inches(1.6), Inches(0.035),
                       fill=self.colors["accent"])

        milestones = body if body else []
        n = min(len(milestones), 6) if milestones else 0
        if n == 0:
            return
        step_x = (SLIDE_W - m * 2 - Inches(1.6)) / max(n - 1, 1)

        for i, ms in enumerate(milestones[:n]):
            cx = m + Inches(0.8) + step_x * i
            node_r = Inches(0.16)
            clr = self.colors["accent"] if i == n - 1 else self.colors["accent_secondary"]

            self._add_circle(slide, cx - node_r, tl_y - node_r, node_r * 2, fill=clr)
            is_above = i % 2 == 0
            label_y = tl_y - Inches(0.9) if is_above else tl_y + Inches(0.35)
            self._add_text(slide, str(ms), cx - Inches(1.1), label_y,
                           Inches(2.2), Inches(0.5),
                           size=Pt(9), color=ink if is_above else ink_subtle,
                           align="center", bold=is_above)

    # ── Rich layout: KPI big-number cards ──

    def _render_kpi_cards(self, slide, s, spec, ct, ink, ink_subtle, surface, is_dark, page_num):
        """3-4 large KPI number cards for data emphasis.

        Card height adapts to remaining slide space so cards never overflow
        the bottom margin, even with airy spacing or large header blocks.
        """
        headline = str(s.get("headline", ""))
        body = s.get("body") or []
        m = self.margin

        body_y = self._draw_header(slide, headline, "", spec, ct, ink, ink_subtle, m)
        cards_data = body if body else []
        n = min(len(cards_data), 4)
        if n == 0:
            return
        gap = self._gap(spec, 0.3)
        card_w = (SLIDE_W - m * 2 - gap * (n - 1)) / n
        card_y = body_y + self._gap(spec, 0.2)

        # ── Adaptive height: never overflow the bottom margin ──
        target_h = Inches(2.4)
        available_h = SLIDE_H - card_y - m
        card_h = min(available_h, target_h)
        if card_h < Inches(1.2):
            # Too little space — fall back to a compact single-row treatment
            card_h = Inches(1.2)

        # Text positions scale proportionally with actual card height
        number_area_h = card_h * 0.55
        label_area_h = card_h * 0.28
        number_y = card_y + card_h * 0.12
        label_y = number_y + number_area_h
        accent_colors = [self.colors["accent"], self.colors["accent_secondary"],
                         self.colors["green"], "#8B5CF6"]
        number_size = Pt(max(18, int(34 * card_h / Inches(2.4))))
        label_size = Pt(max(8, int(10 * card_h / Inches(2.4))))

        for i, item in enumerate(cards_data[:4]):
            cx = m + (card_w + gap) * i
            text = str(item)
            parts = text.split("|", 1) if "|" in text else (text, "")
            number = parts[0].strip()
            label = parts[1].strip() if len(parts) > 1 else ""
            clr = accent_colors[i % len(accent_colors)]

            self._add_rect(slide, cx, card_y, card_w, card_h,
                           fill=surface, line=self.colors["line"], radius=self._slide_corner)
            self._add_rect(slide, cx, card_y, card_w, Inches(0.05), fill=clr)
            self._add_text(slide, number, cx + Inches(0.2), number_y,
                           card_w - Inches(0.4), number_area_h,
                           size=number_size, color=clr, bold=True, align="center")
            if label:
                self._add_text(slide, label, cx + Inches(0.2), label_y,
                               card_w - Inches(0.4), label_area_h,
                               size=label_size, color=ink_subtle, align="center")

    # ── Rich layout: Comparison Table ──

    def _render_comparison_table(self, slide, s, spec, ct, ink, ink_subtle, surface, is_dark, page_num):
        """Side-by-side comparison table."""
        headline = str(s.get("headline", ""))
        body = s.get("body") or []
        m = self.margin

        table_y = self._draw_header(slide, headline, "", spec, ct, ink, ink_subtle, m)

        comparison = s.get("chart_spec", {})
        columns = comparison.get("categories", ["维度", "方案A", "方案B"])
        n_cols = min(len(columns), 5)
        n_rows = min(len(body) if body else 4, 8)
        table_w = SLIDE_W - m * 2
        col_w = []
        if n_cols >= 3:
            col_w = [table_w * 0.25, (table_w * 0.75) / (n_cols - 1), (table_w * 0.75) / (n_cols - 1)]
            if n_cols > 3:
                col_w += [(table_w * 0.75) / (n_cols - 1)] * (n_cols - 3)
        else:
            col_w = [table_w / n_cols] * n_cols

        row_h = Inches(0.48); header_h = Inches(0.45)

        for c, col_name in enumerate(columns[:n_cols]):
            cx = m + sum(col_w[:c]) if c > 0 else m
            bg = self.colors["dark_bg"] if c == 0 else self.colors["accent_soft"]
            tc = self.colors["on_dark"] if c == 0 else ink
            self._add_rect(slide, cx, table_y, col_w[c], header_h, fill=bg)
            self._add_text(slide, str(col_name), cx + Inches(0.15), table_y + Inches(0.06),
                           col_w[c] - Inches(0.3), header_h - Inches(0.12),
                           size=Pt(11), color=tc, bold=True)

        for r in range(n_rows):
            ry = table_y + header_h + row_h * r
            row_bg = surface if r % 2 == 0 else self.colors["bg"]
            for c in range(n_cols):
                cx = m + sum(col_w[:c]) if c > 0 else m
                self._add_rect(slide, cx, ry, col_w[c], row_h, fill=row_bg, line=self.colors["line"])
                cell_text = ""
                # Split body item by pipe for multi-column data
                parts = str(body[r]).split("|") if r < len(body) else []
                if c == 0 and r < len(body):
                    # First column: use first segment (row label) or full text
                    cell_text = parts[0].strip() if parts else str(body[r])
                elif c < len(parts):
                    cell_text = parts[c].strip()
                self._add_text(slide, cell_text, cx + Inches(0.15), ry + Inches(0.06),
                               col_w[c] - Inches(0.3), row_h - Inches(0.12),
                               size=Pt(9), color=ink if c == 0 else ink_subtle, bold=(c == 0))

    # ── Rich layout: Process Flow ──

    def _render_process_flow(self, slide, s, spec, ct, ink, ink_subtle, surface, is_dark, page_num):
        """Horizontal process flow with arrow-connected nodes."""
        headline = str(s.get("headline", ""))
        body = s.get("body") or []
        m = self.margin

        flow_y = self._draw_header(slide, headline, "", spec, ct, ink, ink_subtle, m) + self._gap(spec, 0.3)
        colors = [self.colors["accent"], self.colors["accent_secondary"],
                  self.colors["green"], "#8B5CF6", self.colors["accent"]]

        steps_list = body if body else []
        n = min(len(steps_list), 5)
        if n == 0:
            return

        total_w = SLIDE_W - m * 2
        arrow_gap = Inches(0.3)
        node_w = (total_w - arrow_gap * (n - 1)) / n
        node_h = self._gap(spec, 1.6)
        start_x = m

        for i, step in enumerate(steps_list[:5]):
            sx = start_x + (node_w + arrow_gap) * i
            clr = colors[i % len(colors)]

            self._add_rect(slide, sx, flow_y, node_w, node_h,
                           fill=surface, line=clr, radius=self._slide_corner, name="panel")
            self._add_rect(slide, sx, flow_y, node_w, Inches(0.04), fill=clr, name="card-bar")
            # Number circle
            num_r = Inches(0.2)
            self._add_circle(slide, sx + Inches(0.12), flow_y + Inches(0.12),
                             num_r * 2, fill=clr, name="step-num")
            self._add_text(slide, str(i + 1), sx + Inches(0.12), flow_y + Inches(0.12),
                           num_r * 2, num_r * 2,
                           size=Pt(9), color=self.colors["on_dark"], bold=True, align="center")
            # Step text
            self._add_text(slide, str(step), sx + Inches(0.2), flow_y + Inches(0.55),
                           node_w - Inches(0.4), node_h - Inches(0.7),
                           size=Pt(10), color=ink)
            # Arrow
            if i < n - 1:
                ax = sx + node_w
                ay = flow_y + node_h / 2
                a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                           ax, ay - Inches(0.06),
                                           arrow_gap, Inches(0.12))
                a.fill.solid(); a.fill.fore_color.rgb = _hex_to_rgb(clr)
                a.line.fill.background()

    # ── Rich layout: Quote Callout ──

    def _render_quote_callout(self, slide, s, spec, ct, ink, ink_subtle, surface, is_dark, page_num):
        """Large quote or key statement with decorative elements."""
        headline = str(s.get("headline", ""))
        body = s.get("body") or []
        m = Inches(2.0)

        # Large decorative quote mark
        self._add_text(slide, "“", Inches(0.8), Inches(0.4), Inches(1.5), Inches(1.6),
                       size=Pt(80), color=self.colors["accent"], bold=True)

        # Main quote — larger, centered
        quote_w = SLIDE_W - m * 2
        self._add_text(slide, headline, m, Inches(1.5), quote_w, Inches(2.0),
                       size=Pt(min(34, self.title_max)), color=ink, bold=True)

        # Accent line below quote
        self._add_rect(slide, m, Inches(3.6), Inches(4), Inches(0.045),
                       fill=self.colors["accent"])

        # Attribution / context — smaller, below accent line
        if body:
            attr_y = Inches(3.9)
            for item in body:
                self._add_text(slide, str(item), m, attr_y,
                               quote_w, Inches(0.4),
                               size=Pt(13), color=ink_subtle)
                attr_y += Inches(0.35)

    # ── Rich layout: Agenda / TOC ──

    def _render_agenda(self, slide, s, spec, ct, ink, ink_subtle, surface, is_dark, page_num):
        """Agenda / table of contents with numbered sections."""
        headline = str(s.get("headline", "目录"))
        body = s.get("body") or []
        title_size = spec.get("title_size", 28)
        m = self.margin

        self._add_text(slide, headline, m, Inches(0.5), SLIDE_W - m * 2, Inches(0.6),
                       size=Pt(title_size), color=ink, bold=True)
        self._add_rect(slide, m, Inches(1.15), Inches(2.5), Inches(0.04),
                       fill=self.colors["accent"])

        items = body if body else []
        n = min(len(items), 6)
        if n == 0:
            return
        item_h = (SLIDE_H - Inches(2.0)) / n
        for i, item in enumerate(items[:6]):
            iy = Inches(1.5) + item_h * i
            # Number
            num_r = Inches(0.28)
            self._add_circle(slide, m + Inches(0.1), iy + item_h * 0.3,
                             num_r * 2, fill=self.colors["accent"])
            self._add_text(slide, str(i + 1), m + Inches(0.1), iy + item_h * 0.3,
                           num_r * 2, num_r * 2,
                           size=Pt(14), color=self.colors["on_dark"], bold=True, align="center")
            # Separator line
            if i < n - 1:
                self._add_rect(slide, m + Inches(0.1), iy + item_h,
                               Inches(0.015), item_h * 0.5, fill=self.colors["line"])
            # Text
            self._add_text(slide, str(item), m + Inches(0.9), iy + item_h * 0.1,
                           SLIDE_W - m * 2 - Inches(1.0), item_h - Inches(0.2),
                           size=Pt(self.body_base + 1), color=ink)

    # ── Rich layout: SWOT Quadrant ──

    def _render_swot(self, slide, s, spec, ct, ink, ink_subtle, surface, is_dark, page_num):
        """SWOT analysis 2x2 quadrant."""
        headline = str(s.get("headline", ""))
        body = s.get("body") or []
        title_size = spec.get("title_size", 22)
        m = self.margin

        self._add_text(slide, headline, m, Inches(0.15), SLIDE_W - m * 2, Inches(0.4),
                       size=Pt(title_size), color=ink, bold=True)

        qx = Inches(0.8); qy = Inches(0.9)
        qw = SLIDE_W - Inches(1.6); qh = SLIDE_H - Inches(1.3)
        hw = qw / 2; hh = qh / 2

        quads = [
            ("S 优势", 0, 0, self.colors["green"]),
            ("W 劣势", 1, 0, self.colors["accent_secondary"]),
            ("O 机会", 0, 1, self.colors["accent"]),
            ("T 威胁", 1, 1, "#EF4444"),
        ]

        for label, col, row, clr in quads:
            qx_pos = qx + hw * col; qy_pos = qy + hh * row
            self._add_rect(slide, qx_pos, qy_pos, hw, hh,
                           fill=surface if (col + row) % 2 == 0 else self.colors["bg"],
                           line=self.colors["line"])
            self._add_rect(slide, qx_pos, qy_pos, hw, Inches(0.04), fill=clr)
            self._add_text(slide, label, qx_pos + Inches(0.2), qy_pos + Inches(0.08),
                           Inches(1.5), Inches(0.35), size=Pt(14), color=clr, bold=True)
            ci = col + row * 2
            if ci < len(body):
                self._add_text(slide, str(body[ci]),
                               qx_pos + Inches(0.25), qy_pos + Inches(0.5),
                               hw - Inches(0.5), hh - Inches(0.6),
                               size=Pt(8), color=ink_subtle)

    # ── Rich layout: Funnel ──

    def _render_funnel(self, slide, s, spec, ct, ink, ink_subtle, surface, is_dark, page_num):
        """Funnel diagram for conversion / pipeline."""
        headline = str(s.get("headline", ""))
        body = s.get("body") or []
        title_size = spec.get("title_size", 22)
        m = self.margin

        self._add_text(slide, headline, m, Inches(0.2), SLIDE_W - m * 2, Inches(0.4),
                       size=Pt(title_size), color=ink, bold=True)

        stages = body if body else []
        n = min(len(stages), 5)
        if n == 0:
            return
        max_w = Inches(8.0); min_w = Inches(2.0)
        stage_h = Inches(0.95)
        total_h = stage_h * n + Inches(0.1) * (n - 1)
        start_y = (SLIDE_H - total_h) / 2 + Inches(0.3)
        center_x = SLIDE_W / 2

        for i, stage in enumerate(stages[:5]):
            ratio = 1.0 - (i / max(n - 1, 1)) * 0.8
            fw = max(min_w, max_w * ratio)
            fx = center_x - fw / 2
            fy = start_y + (stage_h + Inches(0.1)) * i
            clr = [self.colors["accent"], self.colors["accent_secondary"],
                   self.colors["green"], "#8B5CF6", self.colors["accent"]][i % 5]

            shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, fx, fy, fw, stage_h)
            shape.fill.solid()
            shape.fill.fore_color.rgb = _hex_to_rgb(clr)
            shape.line.fill.background()

            txt = str(stage)
            self._add_text(slide, txt, fx + Inches(0.5), fy + Inches(0.15),
                           fw - Inches(1.0), stage_h - Inches(0.3),
                           size=Pt(10), color=self.colors["on_dark"], bold=True, align="center")

    # ── Rich layout: 2x2 Matrix ──

    def _render_matrix_2x2(self, slide, s, spec, ct, ink, ink_subtle, surface, is_dark, page_num):
        """2x2 prioritization matrix."""
        headline = str(s.get("headline", ""))
        body = s.get("body") or []
        title_size = spec.get("title_size", 22)
        m = Inches(0.6)

        self._add_text(slide, headline, m, Inches(0.15), SLIDE_W - m * 2, Inches(0.35),
                       size=Pt(title_size), color=ink, bold=True)

        mx = Inches(1.2); my = Inches(0.9)
        mw = Inches(10.2); mh = Inches(5.5)
        mid_x = mx + mw / 2; mid_y = my + mh / 2

        q_data = [
            (mid_x, my, "高价值·易实现", self.colors["green"]),
            (mx, my, "高价值·难实现", self.colors["accent"]),
            (mid_x, mid_y, "低价值·易实现", self.colors["accent_secondary"]),
            (mx, mid_y, "低价值·难实现", "#EF4444"),
        ]

        for (qx1, qy1, label, clr), i in zip(q_data, range(4)):
            qw_q = mid_x - mx if qx1 == mx else mw / 2
            qh_q = mid_y - my if qy1 == my else mh / 2
            self._add_rect(slide, qx1, qy1, qw_q, qh_q,
                           fill=surface if i % 2 == 0 else self.colors["bg"],
                           line=self.colors["line"])
            self._add_text(slide, label, qx1 + Inches(0.15), qy1 + Inches(0.06),
                           Inches(2.5), Inches(0.3), size=Pt(10), color=clr, bold=True)
            if i < len(body):
                self._add_text(slide, str(body[i]), qx1 + Inches(0.2), qy1 + Inches(0.4),
                               qw_q - Inches(0.4), qh_q - Inches(0.5),
                               size=Pt(8), color=ink_subtle)

        # Axes
        self._add_rect(slide, mid_x, my, Inches(0.018), mh, fill=self.colors["ink"])
        self._add_rect(slide, mx, mid_y, mw, Inches(0.018), fill=self.colors["ink"])
        self._add_text(slide, "价值 →", mid_x + Inches(0.1), my + mh + Inches(0.05),
                       Inches(1.5), Inches(0.25), size=Pt(8), color=ink_subtle)
        self._add_text(slide, "可行性 →", mx + mw - Inches(1.2), mid_y - Inches(0.3),
                       Inches(1.2), Inches(0.25), size=Pt(8), color=ink_subtle, align="right")

    # ── Rich layout: Closing CTA ──

    def _render_closing_cta(self, slide, s, spec, ct, ink, ink_subtle, surface, is_dark, page_num):
        """Closing page with strong call-to-action — action items + contact."""
        headline = str(s.get("headline", ""))
        body = s.get("body") or []
        title_size = spec.get("title_size", 32)
        m = self.margin

        # Strong headline
        self._add_text(slide, headline, m, self._gap(spec, 0.25), SLIDE_W - m * 2, self._gap(spec, 0.8),
                       size=Pt(min(title_size, 38)), color=ink, bold=True)

        # Accent bar
        self._add_rect(slide, m, self._gap(spec, 1.2), Inches(3.5), Inches(0.045),
                       fill=self.colors["accent"])

        # Action items
        action_y = self._gap(spec, 2.4)
        for i, item in enumerate(body[:3]):
            num_r = Inches(0.25)
            self._add_circle(slide, m + Inches(0.1), action_y + Inches(0.1),
                             num_r * 2, fill=self.colors["accent"] if i == 0 else self.colors["accent_secondary"])
            self._add_text(slide, str(i + 1), m + Inches(0.1), action_y + Inches(0.1),
                           num_r * 2, num_r * 2,
                           size=Pt(14), color=self.colors["on_dark"], bold=True, align="center")
            self._add_text(slide, str(item), m + Inches(0.9), action_y,
                           SLIDE_W - m * 2 - Inches(1.2), self._gap(spec, 0.6),
                           size=Pt(14), color=ink)
            action_y += self._gap(spec, 0.85)

        # Bottom contact
        contact_y = SLIDE_H - self._gap(spec, 1.2)
        self._add_rect(slide, m, contact_y, SLIDE_W - m * 2, Inches(0.015),
                       fill=self.colors["line"])
        contact_text = str(s.get("speaker_note", "")) or "联系人/团队信息"
        self._add_text(slide, contact_text, m, contact_y + Inches(0.1),
                       SLIDE_W - m * 2, self._gap(spec, 0.4),
                       size=Pt(10), color=ink_subtle, align="center")

    # ── Chart rendering ──

    def _render_chart(self, slide, s, spec, ct, ink):
        chart_spec = s.get("chart_spec", {})
        chart_type = str(chart_spec.get("type", "bar")).lower()
        chart_title = str(chart_spec.get("title", ""))
        categories = chart_spec.get("categories", [])
        values = chart_spec.get("values", [])
        m = self.margin

        if not categories or not values:
            return

        chart_data = self._make_chart_data(categories, values, chart_type)

        structure = resolve_structure(spec)
        proportions = spec.get("proportions", {})
        focal = str((s.get("visual_treatment") or {}).get("focal_element", "body_block"))
        is_chart_focal = focal in ("chart", "data_number")

        # Chart positioning strategies by layout
        if structure in ("title_split", "comparison_split"):
            # In the visual panel area
            visual_w = (SLIDE_W - m * 2) * proportions.get("visual", 0.4)
            cx, cy, cw, ch = (
                SLIDE_W - m - visual_w,
                Inches(1.5),
                visual_w - Inches(0.2),
                SLIDE_H - Inches(1.5) - m,
            )
        elif structure in ("hero_cover", "centered"):
            cx, cy, cw, ch = (Inches(2.5), Inches(3.5), SLIDE_W - Inches(5), Inches(3.0))
        elif structure in ("grid", "vertical_stack"):
            cx, cy, cw = (m, Inches(2.5) if is_chart_focal else Inches(4.0), SLIDE_W - m * 2)
            ch = Inches(3.5) if is_chart_focal else Inches(2.5)
        elif is_chart_focal:
            cx, cy, cw, ch = (m + Inches(0.5), Inches(1.8), SLIDE_W - m * 2 - Inches(1), SLIDE_H - Inches(1.8) - m)
        else:
            cx, cy, cw, ch = (SLIDE_W - m - Inches(5.5), Inches(1.8), Inches(5.5), SLIDE_H - Inches(1.8) - m)

        chart_frame = slide.shapes.add_chart(
            self._chart_type_enum(chart_type),
            cx, cy, cw, ch,
            chart_data
        )

        chart = chart_frame.chart
        chart.has_legend = chart_type in ("pie", "doughnut", "bar", "line")
        chart.has_title = bool(chart_title)
        if chart_title:
            chart.chart_title.text_frame.paragraphs[0].text = chart_title
            chart.chart_title.text_frame.paragraphs[0].font.size = Pt(max(10, self.body_base - 2))

        # Style from design concept
        try:
            if self.concept.decoration_level == "minimal":
                chart.style = 2
            else:
                chart.style = 3
        except Exception:
            pass

        # Apply DesignConcept colors to chart series
        self._apply_chart_colors(chart, chart_type)

        # Data labels for bar/line charts
        if chart_type in ("bar", "line") and chart.series:
            try:
                from pptx.enum.chart import XL_LABEL_POSITION
                for series in chart.series:
                    series.has_data_labels = True
                    series.data_labels.font.size = Pt(8)
                    series.data_labels.font.color.rgb = _hex_to_rgb(self.colors["muted"])
                    series.data_labels.number_format = '0.#'
                    series.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
            except Exception:
                pass

        # Clean gridlines
        try:
            if chart_type in ("bar", "line"):
                from pptx.dml.color import RGBColor
                grid_color = _hex_to_rgb("E5E7EB")
                for axis in (chart.category_axis, getattr(chart, 'value_axis', None)):
                    if axis and getattr(axis, 'has_major_gridlines', False):
                        try:
                            axis.major_gridlines.format.line.color.rgb = grid_color
                            axis.major_gridlines.format.line.width = Pt(0.5)
                        except Exception:
                            pass
        except Exception:
            pass

        # Chart title styling
        if chart_title:
            try:
                chart.chart_title.text_frame.paragraphs[0].font.name = self.font
                chart.chart_title.text_frame.paragraphs[0].font.color.rgb = _hex_to_rgb(self.colors["ink"])
            except Exception:
                pass

    def _apply_chart_colors(self, chart, chart_type: str) -> None:
        """Apply DesignConcept accent colors to chart series."""
        try:
            from pptx.dml.color import RGBColor

            palette = [
                self.colors["accent"],
                self.colors["accent_secondary"],
                self.colors["green"],
                self.colors.get("negative", "DC2626"),
                self.colors.get("muted", "6B7280"),
                self.colors.get("accent_soft", "93C5FD"),
            ]

            for i, series in enumerate(chart.series):
                color_hex = palette[i % len(palette)]
                rgb = _hex_to_rgb(color_hex)
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = rgb

                # For line charts, also set the line color
                if chart_type == "line":
                    series.format.line.color.rgb = rgb
                    series.format.line.width = Pt(2)

                # For pie/doughnut, set point colors individually
                if chart_type in ("pie", "doughnut"):
                    for j, point in enumerate(series.points):
                        pt_color = palette[j % len(palette)]
                        point.format.fill.solid()
                        point.format.fill.fore_color.rgb = _hex_to_rgb(pt_color)
        except Exception:
            pass  # best-effort; chart still renders with defaults if color fails

    # ── Drawing primitives ──

    # EMU conversion constants
    _EMU_PER_INCH = 914400
    _EMU_PER_PT = 12700
    _LINE_HEIGHT_RATIO = 1.6   # line height = font_size × this
    _CJK_CHAR_FACTOR = 0.95    # CJK characters are slightly wider than Latin

    def _estimate_lines(self, text: str, box_w, font_size) -> int:
        """Estimate how many lines `text` needs in a box of width `box_w`."""
        font_pt = font_size / self._EMU_PER_PT
        chars_per_line = max(1, int(box_w / self._EMU_PER_INCH * 72 / font_pt * self._CJK_CHAR_FACTOR))
        return max(1, -(-len(str(text)) // chars_per_line))

    def _text_height(self, text: str, box_w, font_size, given_h) -> int:
        """Return actual height (EMU) needed to fit `text` without truncation."""
        lines = self._estimate_lines(text, box_w, font_size)
        font_pt = font_size / self._EMU_PER_PT
        line_h_emu = int(font_pt / 72 * self._LINE_HEIGHT_RATIO * self._EMU_PER_INCH)
        return max(given_h, lines * line_h_emu)

    def _set_bg(self, slide, color: str) -> None:
        """Set slide background fill."""
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = _hex_to_rgb(color)

    def _add_rect(self, slide, x, y, w, h, fill=None, line=None, radius=None, name=None):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, x, y, w, h)
        if name:
            shape.name = name
        if fill:
            shape.fill.solid()
            shape.fill.fore_color.rgb = _hex_to_rgb(fill)
        else:
            shape.fill.background()
        if line:
            shape.line.color.rgb = _hex_to_rgb(line)
            shape.line.width = Pt(0.5)
        else:
            shape.line.fill.background()
        return shape

    def _add_circle(self, slide, x, y, d, fill=None, name=None):
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
        if name:
            shape.name = name
        if fill:
            shape.fill.solid()
            shape.fill.fore_color.rgb = _hex_to_rgb(fill)
        else:
            shape.fill.background()
        shape.line.fill.background()
        return shape

    def _add_image(self, slide, path, x, y, w, h):
        """Add an image from a local file path. Returns shape or None on failure."""
        try:
            return slide.shapes.add_picture(str(path), x, y, w, h)
        except Exception:
            return None

    def _add_text(self, slide, text, x, y, w, h, size=None, color=None, bold=False, align="left", valign="middle", font=None):
        """Add a text box. Auto-expands height to fit — no text truncation."""
        font_size = size or Pt(self.body_base)
        text_str = str(text)
        actual_h = self._text_height(text_str, w, font_size, int(h))

        txBox = slide.shapes.add_textbox(x, y, w, actual_h)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.clear()
        run = p.add_run()
        run.text = text_str
        run.font.size = font_size
        run.font.color.rgb = _hex_to_rgb(color or self.colors["ink"])
        run.font.bold = bold
        run.font.name = font or self.font
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
        return actual_h

    def _add_body_text(self, slide, body, x, y, w, h, size=None, color=None):
        """Add multi-line body text. Delegates to _add_text for each item."""
        font_size = size or Pt(self.body_base)
        font_color = color or self.colors["muted"]

        # Calculate total height from individual item heights
        total_h = 0
        for item in body:
            total_h += int(self._text_height(str(item), w, font_size, 0))
            total_h += Pt(5)  # paragraph spacing

        actual_h = max(int(h), total_h)

        txBox = slide.shapes.add_textbox(x, y, w, actual_h)
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, item in enumerate(body):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.clear()
            run = p.add_run()
            run.text = str(item)
            run.font.size = font_size
            run.font.color.rgb = _hex_to_rgb(font_color)
            run.font.name = self.font
            p.space_after = Pt(4)
            p.space_before = Pt(1)
        return actual_h

    def _render_body_cards(self, slide, body, x, y, w, h, cols, body_size, ink, ink_subtle, surface):
        """Render body items as a card grid with subtle shadow depth."""
        gap = Inches(0.2)
        n = min(len(body), cols * 2)
        if n == 0 or cols < 1:
            return
        card_w = (w - gap * (cols - 1)) / cols
        rows = (n + cols - 1) // cols
        card_h = (h - gap * (rows - 1)) / rows
        accent_colors = [self.colors["accent"], self.colors["accent_secondary"],
                         self.colors["green"], self.colors["accent"],
                         self.colors["accent_secondary"], self.colors["green"]]

        for i in range(n):
            c = i % cols
            r = i // cols
            cx = x + c * (card_w + gap)
            cy = y + r * (card_h + gap)
            clr = accent_colors[i % len(accent_colors)]

            # Shadow: darker rectangle offset slightly
            shadow_offset = Inches(0.03)
            self._add_rect(slide, cx + shadow_offset, cy + shadow_offset, card_w, card_h,
                           fill="#E0E0E0" if not self.colors["bg"].startswith("0") else "#1A1A2E",
                           radius=self._slide_corner)

            # Card background
            self._add_rect(slide, cx, cy, card_w, card_h,
                           fill=surface, line=self.colors["line"], radius=self._slide_corner)
            # Colored left edge
            self._add_rect(slide, cx, cy, Inches(0.06), card_h, fill=clr)
            # Card content
            self._add_text(slide, str(body[i]), cx + Inches(0.35), cy + Inches(0.15),
                           card_w - Inches(0.65), card_h - Inches(0.3),
                           size=Pt(body_size), color=ink)

    def _make_chart_data(self, categories, values, chart_type):
        """Build a CategoryChartData object for python-pptx."""
        from pptx.chart.data import CategoryChartData

        chart_data = CategoryChartData()
        chart_data.categories = [str(c) for c in categories]

        if values and isinstance(values[0], dict):
            # Multi-series with names
            for i, series in enumerate(values):
                name = str(series.get("name", f"系列{i+1}"))
                data = series.get("data", series.get("values", []))
                if not isinstance(data, (list, tuple)):
                    data = [data]
                chart_data.add_series(name, [_safe_float(v) for v in data])
        elif values and isinstance(values[0], (list, tuple)):
            # Multi-series without names: [[1,2,3], [4,5,6]]
            for i, series in enumerate(values):
                chart_data.add_series(f"系列{i+1}", [_safe_float(v) for v in series])
        else:
            # Single series
            chart_data.add_series("", [_safe_float(v) for v in values])

        return chart_data

    def _chart_type_enum(self, chart_type: str):
        from pptx.enum.chart import XL_CHART_TYPE
        return {
            "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "line": XL_CHART_TYPE.LINE_MARKERS,
            "pie": XL_CHART_TYPE.PIE,
            "doughnut": XL_CHART_TYPE.DOUGHNUT,
            "radar": XL_CHART_TYPE.RADAR,
            "scatter": XL_CHART_TYPE.XY_SCATTER,
        }.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)

    # ── Accent treatments ──

    # ── Layout helpers ──

    def _resolve_margin(self) -> Inches:
        spacing = self.concept.spacing_mood
        mult = self.concept.margin_multiplier
        base = {"airy": 0.9, "normal": 0.6, "compact": 0.45}.get(spacing, 0.6)
        return Inches(base * mult)


# ── Color utility functions ──

def _hex_to_rgb(hex_color: str) -> RGBColor:
    h = str(hex_color).lstrip("#").upper()
    if len(h) == 6:
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    return RGBColor(0x11, 0x18, 0x27)


def _lighten_hex(hex_color: str, amount: float) -> str:
    h = hex_color.lstrip("#")
    if len(h) != 6:
        return "FFFFFF"
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    r = min(255, int(r + (255 - r) * amount))
    g = min(255, int(g + (255 - g) * amount))
    b = min(255, int(b + (255 - b) * amount))
    return f"{r:02X}{g:02X}{b:02X}"


def _muted_hex(ink_hex: str) -> str:
    h = ink_hex.lstrip("#")
    if len(h) != 6:
        return "6B7280"
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    r2 = int(r * 0.55 + 255 * 0.45)
    g2 = int(g * 0.55 + 255 * 0.45)
    b2 = int(b * 0.55 + 255 * 0.45)
    return f"{r2:02X}{g2:02X}{b2:02X}"


def _line_hex(bg_hex: str) -> str:
    h = bg_hex.lstrip("#")
    if len(h) != 6:
        return "E5E7EB"
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    avg = (r + g + b) / 3
    factor = 0.85 if avg > 128 else 1.15
    return f"{min(255, int(r * factor)):02X}{min(255, int(g * factor)):02X}{min(255, int(b * factor)):02X}"


def _panel_hex(bg_hex: str) -> str:
    """Panel color: white for light backgrounds, slightly lighter for dark."""
    h = bg_hex.lstrip("#")
    if len(h) != 6:
        return "FFFFFF"
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    avg = (r + g + b) / 3
    if avg > 200:
        return "FFFFFF"  # already very light, use pure white panels
    if avg > 128:
        # Light bg: panel is noticeably lighter
        return f"{min(255, int(r * 1.08)):02X}{min(255, int(g * 1.08)):02X}{min(255, int(b * 1.08)):02X}"
    # Dark bg: panel is slightly lighter than bg
    return f"{min(255, int(r * 1.25)):02X}{min(255, int(g * 1.25)):02X}{min(255, int(b * 1.25)):02X}"


def _soften_hex(hex_color: str) -> str:
    h = hex_color.lstrip("#")
    if len(h) != 6:
        return "DBEAFE"
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    r2 = int(r * 0.18 + 255 * 0.82)
    g2 = int(g * 0.18 + 255 * 0.82)
    b2 = int(b * 0.18 + 255 * 0.82)
    return f"{r2:02X}{g2:02X}{b2:02X}"


def _safe_float(value: Any) -> float:
    """Convert value to float, returning 0.0 on failure."""
    try:
        return float(value)
    except (TypeError, ValueError):
        return 0.0
